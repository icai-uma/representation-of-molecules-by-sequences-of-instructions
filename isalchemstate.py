"""

Class to model an IsalChem state:
Lopez-Rubio, Ezequiel (2024). Instruction set and language for chemical nomenclature
https://doi.org/10.26434/chemrxiv-2024-5b4dn

Coded by Ezequiel Lopez-Rubio, November 2024.

"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
except ImportError:
    print("The pptx library is not available")

try:
    import networkx as nx
except ImportError:
    print("The networkx library is not available")

from datastructures import Graph, DoublyLinkedCircularList


class IsalChemState:
    """ Class to model an IsalChem state """
    def __init__(self):
        graph = Graph()
        graph.add_node(1, "H")
        graph.add_node(2, "H")
        graph.add_edge(1, 2)

        dll = DoublyLinkedCircularList()
        dll.insert(graph.get_node(1))
        dll.insert_after(graph.get_node(1), graph.get_node(2))

        self.graph = graph
        self.dll = dll
        self.primary_ptr = 1
        self.secondary_ptr = 2

    def __repr__(self):
        graph_repr = self.graph.__repr__()
        dll_repr = self.dll.__repr__()
        return f"IsalChemState(graph={graph_repr}, dll={dll_repr}, primary_ptr={self.primary_ptr}, secondary_ptr={self.secondary_ptr})"

    def generate_pptx(self, filename: str = "graph.pptx"):
        """ Generate a PowerPoint presentation file that represents this IsalChem state """

        # Colors for nodes
        node_colors = {"H": RGBColor(91, 155, 213),  # Light blue color
                       "C": RGBColor(58, 198, 95),  # Green color
                       "N": RGBColor(128, 128, 128),  # Gray color
                       "B": RGBColor(232, 192, 24),  # Orange color
                       "O": RGBColor(224, 0, 0),  # Red color
                       "I": RGBColor(224, 224, 0),  # Yellow color
                       "F": RGBColor(200, 162, 200),  # Lilac color
                       "Cl": RGBColor(31, 225, 207),  # Cyan color
                       "Br": RGBColor(207, 49, 139),  # Magenta color
                       }
        # Create a PowerPoint presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Ensure the slide is blank by removing any placeholder shapes if they exist
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = slide.shapes._spTree.remove(shape._element)

        # Build the NetworkX graph for layout computation
        G = nx.Graph()
        for node in self.graph.nodes.values():
            G.add_node(node.node_id)
        for node in self.graph.nodes.values():
            for neighbor in node.neighbors:
                G.add_edge(node.node_id, neighbor.node_id)

        # Compute positions using spring layout
        pos = nx.spring_layout(G, seed=42)

        # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        margin = Inches(1)

        # Map networkx positions to PowerPoint coordinates
        node_shapes = {}
        for node in self.graph.nodes.values():
            # Scale node positions to fit in slide dimensions
            x = pos[node.node_id][0] * (slide_width - 2 * margin) + slide_width / 2
            y = pos[node.node_id][1] * (slide_height - 2 * margin) + slide_height / 2

            # For the primary pointer,
            # add a square enclosing shape slightly larger than the circle
            if node.node_id == self.primary_ptr:
                square_size = Inches(0.6)  # Larger than circle size
                square = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(x - square_size / 2),
                    int(y - square_size / 2),
                    int(square_size),
                    int(square_size)
                )
                square.fill.background()  # Transparent fill
                square.line.color.rgb = RGBColor(255, 0, 0)  # Red outline
                square.line.width = Pt(2)  # Thicker outline

            # For the next of the primary pointer,
            # add a square enclosing shape slightly larger than the circle
            if node.prev is not None and node.prev.node_id == self.primary_ptr:
                square_size = Inches(0.6)  # Larger than circle size
                square = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(x - square_size / 2),
                    int(y - square_size / 2),
                    int(square_size),
                    int(square_size)
                )
                square.fill.background()  # Transparent fill
                square.line.color.rgb = RGBColor(0, 255, 0)  # Green outline
                square.line.width = Pt(2)  # Thicker outline

            # For the secondary pointer,
            # add a square enclosing shape slightly larger than the circle
            if node.node_id == self.secondary_ptr:
                square_size = Inches(0.7)  # Larger than circle size
                square = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(x - square_size / 2),
                    int(y - square_size / 2),
                    int(square_size),
                    int(square_size)
                )
                square.fill.background()  # Transparent fill
                square.line.color.rgb = RGBColor(0, 0, 255)  # Blue outline
                square.line.width = Pt(2)  # Thicker outline

            # Add a circular shape for the node
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(x - Inches(0.25)),
                int(y - Inches(0.25)),
                int(Inches(0.5)),
                int(Inches(0.5))
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = node_colors[node.data]  # Node fill color
            circle.text = node.data

            # Format the text inside the circle
            text_frame = circle.text_frame
            text_frame.text = node.data
            # For debugging purposes
            circle.text = f"{node.node_id}"
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

            # Store the shape for the node
            node_shapes[node] = circle

        # Draw edges as solid connectors between node shapes for chemical bonds
        for node in self.graph.nodes.values():
            for ndx_neighbor, neighbor in enumerate(node.neighbors):
                if node.node_id < neighbor.node_id:  # Avoid duplicate connectors
                    start_shape = node_shapes[node]
                    end_shape = node_shapes[neighbor]

                    # Create a straight connector
                    connector = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT,
                        0, 0, 0, 0
                    )
                    # Connect to the middle of each circle
                    connector.begin_connect(start_shape, 2)  # 2 = middle connection point
                    connector.end_connect(end_shape, 2)
                    connector.line.color.rgb = RGBColor(0, 0, 0)  # Black color for lines
                    # Add thicker outline for double bonds
                    if node.edge_types[ndx_neighbor] == 2:
                        connector.line.width = Pt(6)  # Thicker outline
                    # Add even thicker outline for triple bonds
                    if node.edge_types[ndx_neighbor] == 3:
                        connector.line.width = Pt(8)  # Thicker outline

        # Draw edges as dashed connectors between node shapes for the list connections
        current = self.dll.head
        if current is not None:
            while True:
                start_shape = node_shapes[current]
                end_shape = node_shapes[current.next]

                # Create a straight connector
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    0, 0, 0, 0
                )
                # Connect to the middle of each circle
                connector.begin_connect(start_shape, 2)  # 2 = middle connection point
                connector.end_connect(end_shape, 2)
                connector.line.color.rgb = RGBColor(0, 0, 0)  # Black color for lines
                connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dashed line style

                current = current.next
                if current == self.dll.head:
                    break

        # Save the presentation
        prs.save(filename)
        print(f"Graph saved to {filename}")

